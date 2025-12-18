import os
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import pandas as pd

FOLDER = '.'
pptx_files = [f for f in os.listdir(FOLDER) if f.lower().endswith('.pptx')]

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

texts = []
for file in pptx_files:
    text = extract_text_from_pptx(os.path.join(FOLDER, file))
    texts.append(text)

if len(texts) > 1:
    vectorizer = TfidfVectorizer().fit_transform(texts)
    similarity_matrix = cosine_similarity(vectorizer)
    similarity_df = pd.DataFrame(similarity_matrix, index=pptx_files, columns=pptx_files)
    similarity_df = (similarity_df * 100).round(2)
    similarity_df.to_excel('pptx_similarity.xlsx')
    similarity_df.to_csv('pptx_similarity.csv')
    print('Similarity comparison table saved as pptx_similarity.xlsx and pptx_similarity.csv')
else:
    print('Not enough PPTX files to compare.')
