import os
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import pandas as pd

# Folder containing PPTX files
FOLDER = 'Ch2'

# Get all pptx files
pptx_files = [f for f in os.listdir(FOLDER) if f.lower().endswith('.pptx')]
print(f"Found {len(pptx_files)} PPTX files: {pptx_files}")

# Function to extract all text from a pptx file
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

# Extract text from all files
texts = []
for file in pptx_files:
    try:
        text = extract_text_from_pptx(os.path.join(FOLDER, file))
        texts.append(text)
        print(f"Extracted {len(text)} chars from {file}")
    except Exception as e:
        print(f"Error extracting from {file}: {e}")
        texts.append("")

print(f"Total texts: {len(texts)}")

if len(texts) > 1 and any(len(t) > 0 for t in texts):
    try:
        vectorizer = TfidfVectorizer().fit_transform(texts)
        similarity_matrix = cosine_similarity(vectorizer)
        similarity_df = pd.DataFrame(similarity_matrix, index=pptx_files, columns=pptx_files)
        similarity_df = (similarity_df * 100).round(2)
        similarity_df.to_excel('pptx_similarity.xlsx')
        similarity_df.to_csv('pptx_similarity.csv')
        print('Similarity comparison table saved as pptx_similarity.xlsx and pptx_similarity.csv')
        
        # Generate detailed report
        with open('similarity_report.txt', 'w', encoding='utf-8') as f:
            f.write("DETAILED SIMILARITY REPORT FOR PPTX FILES\n")
            f.write("=" * 50 + "\n\n")
            f.write(f"Folder: {FOLDER}\n")
            f.write(f"Number of files: {len(pptx_files)}\n\n")
            
            # Summary statistics
            f.write("SUMMARY STATISTICS:\n")
            f.write("-" * 20 + "\n")
            f.write(f"Average similarity: {similarity_df.values.mean():.2f}%\n")
            f.write(f"Max similarity: {similarity_df.values.max():.2f}%\n")
            f.write(f"Min similarity: {similarity_df.values.min():.2f}%\n\n")
            
            # High similarity pairs (>70%)
            f.write("HIGH SIMILARITY PAIRS (>70%):\n")
            f.write("-" * 30 + "\n")
            high_sim_pairs = []
            for i in range(len(pptx_files)):
                for j in range(i+1, len(pptx_files)):
                    sim = similarity_df.iloc[i, j]
                    if sim > 70:
                        high_sim_pairs.append((pptx_files[i], pptx_files[j], sim))
                        f.write(f"{pptx_files[i]} vs {pptx_files[j]}: {sim:.2f}%\n")
            
            if not high_sim_pairs:
                f.write("No pairs with similarity > 70%\n")
            f.write("\n")
            
            # Potential plagiarism concerns (>80%)
            f.write("POTENTIAL PLAGIARISM CONCERNS (>80%):\n")
            f.write("-" * 35 + "\n")
            plagiarism_pairs = [pair for pair in high_sim_pairs if pair[2] > 80]
            for pair in plagiarism_pairs:
                f.write(f"{pair[0]} vs {pair[1]}: {pair[2]:.2f}% - HIGH RISK\n")
            
            if not plagiarism_pairs:
                f.write("No pairs with similarity > 80%\n")
            f.write("\n")
            
            # Deduction points for each file based on max similarity
            f.write("DEDUCTION POINTS (0-5) BASED ON MAX SIMILARITY:\n")
            f.write("-" * 45 + "\n")
            for file in pptx_files:
                # Max similarity with other files
                max_sim = 0
                for other in pptx_files:
                    if other != file:
                        sim = similarity_df.loc[file, other]
                        if sim > max_sim:
                            max_sim = sim
                # Calculate deduction points
                if max_sim < 20:
                    points = 0
                elif max_sim < 40:
                    points = 1
                elif max_sim < 60:
                    points = 2
                elif max_sim < 80:
                    points = 3
                elif max_sim < 90:
                    points = 4
                else:
                    points = 5
                f.write(f"{file}: Max similarity {max_sim:.2f}% -> Deduction: {points} points\n")
            f.write("\n")
            
            # Full similarity matrix
            f.write("FULL SIMILARITY MATRIX (%):\n")
            f.write("-" * 25 + "\n")
            f.write(similarity_df.to_string() + "\n")
        
        print('Detailed report saved as similarity_report.txt')
    except Exception as e:
        print(f"Error in similarity calculation: {e}")
else:
    print('Not enough valid PPTX files to compare.')


